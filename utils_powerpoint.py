import os
import tempfile
from pptx import Presentation  # For handling PowerPoint documents
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_ollama import OllamaEmbeddings
from langchain_chroma import Chroma

# Function to load, split, and embed data from PowerPoint documents into Chroma vector store
def process_powerpoint_documents(powerpoints):
    """
    Process PowerPoint documents through loading, splitting, and embedding.
    Returns vector store instance.
    """
    # Create temporary directory for PowerPoint storage
    with tempfile.TemporaryDirectory() as temp_dir:
        # Save uploaded PowerPoint documents to temp directory
        ppt_paths = []
        for ppt in powerpoints:
            path = os.path.join(temp_dir, ppt.name)
            with open(path, "wb") as f:
                f.write(ppt.getbuffer())
            ppt_paths.append(path)
        
        # Load the documents
        documents = []
        for path in ppt_paths:
            try:
                presentation = Presentation(path)
                text = ""
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text.strip() + "\n"
                if text.strip():  # Only add non-empty text
                    documents.append(text)
                else:
                    print(f"Warning: Document '{path}' contains no readable text.")
            except Exception as e:
                print(f"Error processing PowerPoint file '{path}': {str(e)}")
        
        # Check if there is any valid text to process
        if not documents:
            raise ValueError("No valid text found in the uploaded PowerPoint documents.")
        
        # Split documents into chunks using RecursiveCharacterTextSplitter
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1200,  
            chunk_overlap=150  
        )
        splits = text_splitter.split_text("\n".join(documents))
        
        # Validate splits
        if not splits:
            raise ValueError("No valid chunks generated from the PowerPoint documents.")
        
        # Instantiate the embeddings model
        embeddings = OllamaEmbeddings(model="deepseek-r1:1.5b")
        
        # Create embeddings and vector store
        vector_store = Chroma.from_texts(
            texts=splits,
            embedding=embeddings,
            persist_directory="./chroma_db"
        )
        return vector_store

# Initialize and returns a retriever for the vector store
def get_retriever():
    """Initialize and return the vector store retriever"""
    # Initialize the embedding model
    embeddings = OllamaEmbeddings(model="deepseek-r1:1.5b")
    try:
        # Initialize the vector store
        vector_store = Chroma(
            embedding_function=embeddings,
            persist_directory="./chroma_db"
        )
        # Return the retriever with MMR (Maximum Marginal Relevance) search and k=3
        return vector_store.as_retriever(search_type="mmr", search_kwargs={"k": 3})
    except Exception as e:
        print(f"Error initializing vector store: {e}")
        return None